import os
from datetime import datetime, timezone
from typing import List

import matplotlib.pyplot as plt
import pandas as pd
import seaborn as sns
from matplotlib import gridspec
from pptx import Presentation
from pptx.util import Pt

BASE_DIR = os.path.dirname(__file__)
DATA_DIR = os.path.join(BASE_DIR, "mock_data")
REPORT_DIR = os.path.join(BASE_DIR, "report")
TIMESTAMP_FMT = "%Y-%m-%d %H:%M:%S"

POSITIVE_TERMS = [
    "love",
    "great",
    "fantastic",
    "game changer",
    "intuitive",
    "helped",
    "flawless",
    "super",
    "thanks",
]
NEGATIVE_TERMS = [
    "slow",
    "crash",
    "noisy",
    "issue",
    "problem",
    "confusing",
    "friction",
    "hard",
    "broken",
]


def load_datasets():
    users = pd.read_csv(os.path.join(DATA_DIR, "users.csv"), parse_dates=["signup_date"])
    sessions = pd.read_csv(
        os.path.join(DATA_DIR, "sessions.csv"), parse_dates=["start_time", "end_time"]
    )
    feature_usage = pd.read_csv(
        os.path.join(DATA_DIR, "feature_usage.csv"), parse_dates=["usage_timestamp"]
    )
    feedback = pd.read_csv(os.path.join(DATA_DIR, "feedback.csv"))
    sessions = sessions.merge(
        users[["user_id", "signup_date"]], on="user_id", how="left", validate="many_to_one"
    )
    return users, sessions, feature_usage, feedback


def compute_dau_trend(sessions: pd.DataFrame) -> pd.DataFrame:
    dau = (
        sessions.assign(date=sessions["start_time"].dt.date)
        .groupby("date")["user_id"]
        .nunique()
        .rename("dau")
        .to_frame()
    )
    dau.index = pd.to_datetime(dau.index)
    dau["dau_7d"] = dau["dau"].rolling(window=7).mean()
    return dau


def compute_feature_trend(feature_usage: pd.DataFrame) -> pd.DataFrame:
    feature_usage = feature_usage.copy()
    feature_usage["week"] = feature_usage["usage_timestamp"].dt.to_period("W").dt.start_time
    top_features = feature_usage["feature_name"].value_counts().head(4).index.tolist()
    weekly = (
        feature_usage[feature_usage["feature_name"].isin(top_features)]
        .groupby(["week", "feature_name"])["session_id"]
        .count()
        .rename("events")
        .reset_index()
    )
    return weekly


def compute_retention_heatmap(sessions: pd.DataFrame) -> pd.DataFrame:
    active_sessions = sessions[sessions["start_time"] >= sessions["signup_date"]].copy()
    active_sessions["cohort"] = active_sessions["signup_date"].dt.to_period("M")
    active_sessions["active_month"] = active_sessions["start_time"].dt.to_period("M")
    active_sessions["month_offset"] = (
        active_sessions["active_month"].astype(int) - active_sessions["cohort"].astype(int)
    )
    active_sessions = active_sessions[active_sessions["month_offset"] >= 0]

    cohort_sizes = (
        active_sessions.groupby("cohort")["user_id"].nunique().rename("cohort_size")
    )
    retained = (
        active_sessions.groupby(["cohort", "month_offset"])["user_id"]
        .nunique()
        .unstack(fill_value=0)
        .sort_index()
    )
    retention_rate = retained.divide(cohort_sizes, axis=0)
    retention_rate.columns = [f"M+{int(col)}" for col in retention_rate.columns]
    retention_rate.index = retention_rate.index.astype(str)
    return retention_rate


def classify_feedback(row) -> str:
    rating = row.get("rating", 0)
    comment = str(row.get("comments", "")).lower()
    score = 0
    if rating >= 4:
        score += 1
    elif rating <= 2:
        score -= 1

    score += sum(term in comment for term in POSITIVE_TERMS)
    score -= sum(term in comment for term in NEGATIVE_TERMS)

    if score > 0:
        return "Positive"
    if score < 0:
        return "Negative"
    return "Neutral"


def compute_feedback_trend(feedback: pd.DataFrame, sessions: pd.DataFrame) -> pd.DataFrame:
    feedback = feedback.merge(
        sessions[["session_id", "start_time"]], on="session_id", how="left"
    )
    feedback["month"] = feedback["start_time"].dt.to_period("M").dt.to_timestamp()
    feedback["sentiment"] = feedback.apply(classify_feedback, axis=1)
    summary = (
        feedback.groupby(["month", "sentiment"])['feedback_id']
        .count()
        .rename("count")
        .reset_index()
    )
    return summary


def load_funnel_table() -> pd.DataFrame:
    funnel_path = os.path.join(REPORT_DIR, "funnel_table.csv")
    if os.path.exists(funnel_path):
        funnel = pd.read_csv(funnel_path)
    else:
        raise FileNotFoundError(
            "Funnel data not found. Run generate_funnel.py before generating dashboard."
        )
    return funnel


def render_dashboard(
    dau: pd.DataFrame,
    feature_trend: pd.DataFrame,
    retention: pd.DataFrame,
    funnel: pd.DataFrame,
    feedback_trend: pd.DataFrame,
    output_path: str,
):
    sns.set_theme(style="whitegrid")
    fig = plt.figure(figsize=(16, 9))
    gs = gridspec.GridSpec(2, 3, figure=fig, height_ratios=[1, 1.05], width_ratios=[1.1, 1.1, 0.9])

    # DAU trend
    ax1 = fig.add_subplot(gs[0, 0])
    ax1.plot(dau.index, dau["dau"], label="DAU", color="#5b8def", alpha=0.4)
    ax1.plot(dau.index, dau["dau_7d"], label="7-day avg", color="#1d3b75", linewidth=2)
    ax1.set_title("Daily Active Users")
    ax1.set_ylabel("Users")
    ax1.legend(loc="upper left")

    # Feature usage trend
    ax2 = fig.add_subplot(gs[0, 1])
    if not feature_trend.empty:
        for feature, group in feature_trend.groupby("feature_name"):
            ax2.plot(group["week"], group["events"], label=feature)
        ax2.legend(loc="upper left")
    ax2.set_title("Weekly Feature Usage (Top 4)")
    ax2.set_ylabel("Events")

    # Funnel chart
    ax3 = fig.add_subplot(gs[0, 2])
    bars = ax3.barh(funnel["stage"], funnel["users"], color="#5b8def")
    ax3.invert_yaxis()
    ax3.set_title("Funnel Conversion")
    ax3.set_xlabel("Users")
    for bar, drop in zip(bars, funnel["step_drop_pct"]):
        ax3.text(
            bar.get_width() * 0.99,
            bar.get_y() + bar.get_height() / 2,
            f"{int(bar.get_width()):,}\nDrop {drop:.1f}%",
            va="center",
            ha="right",
            color="white",
            fontweight="bold",
        )

    # Retention heatmap (two columns wide)
    ax4 = fig.add_subplot(gs[1, :2])
    sns.heatmap(
        retention,
        ax=ax4,
        cmap="Blues",
        annot=True,
        fmt=".0%",
        cbar_kws={"label": "Retention"},
    )
    ax4.set_title("Monthly Cohort Retention")
    ax4.set_xlabel("Months Since Signup")
    ax4.set_ylabel("Signup Cohort")

    # Feedback sentiment
    ax5 = fig.add_subplot(gs[1, 2])
    if not feedback_trend.empty:
        pivot = feedback_trend.pivot(index="month", columns="sentiment", values="count").fillna(0)
        pivot = pivot.sort_index()
        sentiments = ["Positive", "Neutral", "Negative"]
        bottom = None
        colors = {"Positive": "#4caf50", "Neutral": "#ffb300", "Negative": "#e53935"}
        for sentiment in sentiments:
            values = pivot.get(sentiment, pd.Series(0, index=pivot.index))
            ax5.bar(pivot.index, values, bottom=bottom, label=sentiment, color=colors[sentiment])
            bottom = values if bottom is None else bottom + values
        ax5.legend(loc="upper left")
    ax5.set_title("Feedback Volume & Sentiment")
    ax5.set_ylabel("Feedback Count")

    fig.suptitle("Product Health Dashboard", fontsize=18, fontweight="bold")
    fig.autofmt_xdate(rotation=30)
    plt.tight_layout(rect=[0, 0, 1, 0.97])

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    fig.savefig(output_path, dpi=150)
    plt.close(fig)


def create_summary_slide(prs: Presentation, title: str, bullet_points: List[str]):
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.text = bullet_points[0]
    for point in bullet_points[1:]:
        p = body.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(18)
    for paragraph in body.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.name = "Calibri"


def build_executive_summary(
    dau: pd.DataFrame,
    feature_trend: pd.DataFrame,
    retention: pd.DataFrame,
    feedback_trend: pd.DataFrame,
    output_path: str,
):
    latest_dau = int(dau["dau"].iloc[-1]) if not dau.empty else 0
    rolling_dau = float(dau["dau_7d"].iloc[-1]) if not dau.empty else 0
    dau_chg = 0.0
    if dau.shape[0] > 30:
        base = dau["dau_7d"].iloc[-31]
        if base:
            dau_chg = (rolling_dau - base) / base

    top_features = (
        feature_trend.groupby("feature_name")["events"].sum().sort_values(ascending=False)
        if not feature_trend.empty
        else pd.Series(dtype=float)
    )
    top_feature_text = ", ".join(
        f"{name} ({int(total):,})" for name, total in top_features.head(3).items()
    )

    avg_month1_retention = retention["M+1"].mean() if "M+1" in retention.columns else 0
    feedback_totals = (
        feedback_trend.groupby("sentiment")["count"].sum()
        if not feedback_trend.empty
        else pd.Series(dtype=int)
    )
    negative_share = 0.0
    if not feedback_totals.empty:
        total_feedback = feedback_totals.sum()
        negative_share = feedback_totals.get("Negative", 0) / total_feedback if total_feedback else 0

    prs = Presentation()
    create_summary_slide(
        prs,
        "Key Insights",
        [
            f"Latest DAU: {latest_dau:,} (7-day avg {rolling_dau:,.0f}, {dau_chg:+.1%} vs. prior 30-day avg)",
            f"Top engaged features: {top_feature_text or 'N/A'}",
            f"Month+1 retention averages {avg_month1_retention:.1%} across cohorts",
        ],
    )
    create_summary_slide(
        prs,
        "Risks & Sprint Agenda",
        [
            f"72% funnel drop from account creation to first transaction remains primary risk",
            f"Negative sentiment share at {negative_share:.1%}; prioritize notification noise fixes",
            "Sprint focus: Guided first-transaction checklist, notification batching, proactive onboarding nudges",
        ],
    )
    prs.save(output_path)


def main():
    os.makedirs(REPORT_DIR, exist_ok=True)

    users, sessions, feature_usage, feedback = load_datasets()
    dau = compute_dau_trend(sessions)
    feature_trend = compute_feature_trend(feature_usage)
    retention = compute_retention_heatmap(sessions)
    funnel = load_funnel_table()
    feedback_trend = compute_feedback_trend(feedback, sessions)

    dashboard_path = os.path.join(REPORT_DIR, "dashboard_overview.png")
    render_dashboard(dau, feature_trend, retention, funnel, feedback_trend, dashboard_path)

    summary_path = os.path.join(REPORT_DIR, "executive_summary.pptx")
    build_executive_summary(dau, feature_trend, retention, feedback_trend, summary_path)

    # Export source data snapshot for reference
    dau_out = os.path.join(REPORT_DIR, "dashboard_dau.csv")
    dau.to_csv(dau_out)

    feature_out = os.path.join(REPORT_DIR, "dashboard_feature_trend.csv")
    feature_trend.to_csv(feature_out, index=False)

    retention_out = os.path.join(REPORT_DIR, "dashboard_retention.csv")
    retention.to_csv(retention_out)

    feedback_out = os.path.join(REPORT_DIR, "dashboard_feedback_trend.csv")
    feedback_trend.to_csv(feedback_out, index=False)

    print("Dashboard assets generated:")
    print("- Overview chart:", dashboard_path)
    print("- Executive summary slides:", summary_path)
    print("- Data extracts saved in report/ for reference")


if __name__ == "__main__":
    main()
